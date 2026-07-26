import streamlit as st
import os
import shutil
import tempfile
import json
import re
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy

# 直接从 skill 文件夹导入（不用 sys.path）
from skill.turing_ppt_editor import (
    set_cover_title,
    set_cover_school,
    copy_external_slide,
    map_theme_to_ch2,
    resolve_out_path
)

# ==================== 1. 配置 DeepSeek 客户端 ====================
client = OpenAI(
    api_key=st.secrets["DEEPSEEK_API_KEY"],
    base_url="https://api.deepseek.com"
)

# ==================== 2. 领域与细分映射 ====================
SUB_OPTIONS = {
    "具身智能": ["部署教学", "场景开发", "仿真训练", "数据采集", "维修"],
    "人工智能": [],
    "低空经济": []
}

# ==================== 3. 文件解析函数（保持不变） ====================
def extract_text_from_file(uploaded_file):
    import os
    import docx
    import pandas as pd
    from pptx import Presentation as PPTXReader
    import PyPDF2
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    text = ""
    try:
        if ext == ".txt":
            text = uploaded_file.getvalue().decode("utf-8")
        elif ext == ".docx":
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in [".xls", ".xlsx"]:
            df = pd.read_excel(uploaded_file)
            text = df.to_string(index=False)
        elif ext == ".pptx":
            prs = PPTXReader(uploaded_file)
            slides_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            text = "\n\n".join(slides_text)
        elif ext == ".pdf":
            reader = PyPDF2.PdfReader(uploaded_file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif ext == ".csv":
            df = pd.read_csv(uploaded_file)
            text = df.to_string(index=False)
        else:
            st.error(f"不支持的文件格式：{ext}")
            return None
    except Exception as e:
        st.error(f"文件读取失败：{e}")
        return None

    file_size = uploaded_file.size
    max_size = 200 * 1024 * 1024
    if file_size > max_size:
        st.info("文件超过200MB，已截取前3000字符作为参考。")
        if len(text) > 3000:
            text = text[:3000] + "...(内容已截断)"
    return text

# ==================== 4. 网页UI ====================
st.title("智能PPT生成器")
st.markdown("填写以下信息，AI将为您自动生成定制化PPT。")

col1, col2 = st.columns(2)

with col1:
    school_name = st.text_input("学校名称", placeholder="请输入学校名称")
    school_level = st.selectbox(
        "学校层次",
        ["985", "211", "双一流", "重点一本", "公办本科", "民办本科", "专科"],
        index=None,
        placeholder="请选择学校层次"
    )

with col2:
    pass

domain = st.selectbox(
    "领域选择",
    list(SUB_OPTIONS.keys()),
    index=None,
    placeholder="请选择领域"
)

sub_domain = None
if domain == "具身智能":
    sub_domain = st.selectbox(
        "细分方向",
        SUB_OPTIONS["具身智能"],
        index=None,
        placeholder="请选择细分方向"
    )

need_nvidia = st.selectbox(
    "是否体现英伟达技术/合作",
    ["是", "否"],
    index=None,
    placeholder="请选择"
)

budget = st.number_input("项目预算（万元）", min_value=0, value=None, step=10, placeholder="请输入预算")
cost = st.number_input("目标成本控制（万元）", min_value=0, value=None, step=10, placeholder="请输入成本")

uploaded_image = st.file_uploader("场地地形图片（可选）", type=["png", "jpg", "jpeg"])
if uploaded_image:
    st.image(uploaded_image, caption="预览上传图片", width=300)

uploaded_files = st.file_uploader(
    "上传参考资料（Word/PPT/Excel/PDF/TXT/CSV，可多选）",
    type=["docx", "pptx", "xls", "xlsx", "pdf", "txt", "csv"],
    accept_multiple_files=True
)

reference_texts = []
if uploaded_files:
    for file in uploaded_files:
        with st.spinner(f"正在读取 {file.name} ..."):
            content = extract_text_from_file(file)
            if content:
                reference_texts.append(f"【文件：{file.name}】\n{content}")

extra_requirements = st.text_area(
    "补充需求（可选）",
    placeholder="请输入补充需求",
    height=100,
    help="您可以在这里输入额外的要求，例如：重点突出技术路线、增加对比图表等。"
)

# ==================== 5. AI 提示构造 ====================
def build_ai_prompt_for_editor(school_name, school_level, domain, sub_domain,
                               budget, cost, need_nvidia, extra_requirements,
                               reference_texts=None):
    nvidia_text = "是" if need_nvidia == "是" else "否"
    refs = "\n".join(reference_texts) if reference_texts else "无"
    prompt = f"""你是一位专业的PPT方案策划专家。请根据以下信息，生成PPT的标题和目录章节名称（第3-6章）。

【用户信息】
- 学校名称：{school_name}
- 学校层次：{school_level if school_level else '未指定'}
- 核心领域：{domain if domain else '未指定'}，细分方向：{sub_domain if sub_domain else '无'}
- 项目预算：{budget if budget else '未指定'}万元，目标成本控制：{cost if cost else '未指定'}万元
- 是否体现英伟达：{nvidia_text}
- 用户补充需求：{extra_requirements if extra_requirements else '无'}

【参考资料】
{refs}

【请严格按照以下格式输出，只输出JSON，不要有其他文字】
{{
  "ppt_title": "这里写PPT的完整标题（建议包含学校名称和领域）",
  "chapters_tail": ["第3章具体标题", "第4章具体标题", "第5章具体标题", "第6章具体标题"]
}}

注意：
1. 第1章固定为"图灵智新介绍"，第2章由系统根据领域自动生成（如"具身智能产业背景及国家政策"），你只需生成第3-6章的标题。
2. 生成的章节标题要具体、专业，结合细分方向（如"数据采集技术路线"）。
3. 标题要体现学校名称和核心领域。
"""
    return prompt

def call_ai_for_editor_content(prompt):
    system_prompt = "你是一个严格的JSON生成器，必须返回合法JSON。"
    try:
        response = client.chat.completions.create(
            model="deepseek-v4-pro",  # 使用最新模型名
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1500
        )
        raw = response.choices[0].message.content
        if raw.startswith("```json"):
            raw = raw[7:-3].strip()
        elif raw.startswith("```"):
            raw = raw[3:-3].strip()
        return json.loads(raw)
    except Exception as e:
        st.error(f"AI解析失败: {e}")
        return None

# ==================== 6. 合并PPT函数（使用学长 copy_external_slide） ====================
def merge_pptx_files(file_paths, output_path,
                     replace_cover_title=None, replace_school_name=None):
    """
    合并多个PPT文件，保留所有内容（图片、背景等）。
    仅修改封面标题和学校名（若提供），其他页面原样拼接。
    """
    if not file_paths:
        raise ValueError("至少需要一个PPT文件")

    # 加载第一个文件作为目标
    target_prs = Presentation(file_paths[0])

    # 依次复制后续文件的所有幻灯片
    for fpath in file_paths[1:]:
        src_prs = Presentation(fpath)
        for slide in src_prs.slides:
            copy_external_slide(target_prs, slide, src_prs)

    # 修改封面（第1张幻灯片）
    if replace_cover_title:
        set_cover_title(target_prs.slides[0], replace_cover_title)
    if replace_school_name:
        set_cover_school(target_prs.slides[0], replace_school_name)

    # 保存
    target_prs.save(output_path)
    return target_prs

def add_image_slide_to_end(prs, image_data):
    """在PPT末尾添加一页带图片的幻灯片（标题：场地地形）"""
    try:
        layout = prs.slide_layouts[6]  # 空白布局
    except IndexError:
        layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # 清空默认形状
    for shp in list(slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # 添加标题
    left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "场地地形"
    p.font.size = Pt(28)
    p.font.bold = True
    # 插入图片
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        tmp.write(image_data)
        tmp_path = tmp.name
    slide.shapes.add_picture(tmp_path, Inches(1), Inches(1.8), width=Inches(8))
    os.unlink(tmp_path)

# ==================== 7. 生成按钮逻辑 ====================
if st.button("生成 PPT"):
    if not school_name.strip():
        st.warning("请输入学校名称")
    elif not domain:
        st.warning("请选择领域")
    elif domain == "具身智能" and not sub_domain:
        st.warning("请选择具身智能的细分方向")
    else:
        with st.spinner("正在生成定制化内容..."):
            # 7.1 调用AI生成标题和章节
            ai_prompt = build_ai_prompt_for_editor(
                school_name, school_level, domain, sub_domain,
                budget, cost, need_nvidia, extra_requirements,
                reference_texts=reference_texts if uploaded_files else None
            )
            ai_result = call_ai_for_editor_content(ai_prompt)
            if ai_result is None:
                st.error("AI生成内容失败，请重试")
                st.stop()

            ppt_title = ai_result.get("ppt_title", f"{school_name} - {domain}建设方案")
            chapters_tail = ai_result.get("chapters_tail", ["技术路线", "实施方案", "预期成果", "保障措施"])

            # 构造完整的6章目录（仅用于AI，实际不修改目录页）
            ch1 = "图灵智新介绍"
            ch2 = map_theme_to_ch2(domain)
            full_chapters = [ch1, ch2] + chapters_tail[:4]
            if len(full_chapters) < 6:
                full_chapters += [f"第{i+1}章" for i in range(6 - len(full_chapters))]
            full_chapters = full_chapters[:6]
            # 注意：这里生成的 full_chapters 没有实际用于修改目录，只是保留以备后用

            # 7.2 确定拼接文件列表
            moban_dir = os.path.join(os.getcwd(), "moban")
            file_list = []

            # 构造文件名前缀
            if domain == "具身智能":
                prefix = f"具身智能{sub_domain}"
            elif domain == "人工智能":
                prefix = "人工智能"
            elif domain == "低空经济":
                prefix = "低空经济"
            else:
                prefix = domain

            # 1. 首页：标题目录 + 第一章
            title_file = os.path.join(moban_dir, f"标题目录{prefix}.pptx")
            if not os.path.exists(title_file):
                st.error(f"找不到首页文件：{title_file}")
                st.stop()
            file_list.append(title_file)

            chapter1_file = os.path.join(moban_dir, "第一章.pptx")
            if not os.path.exists(chapter1_file):
                st.error("找不到第一章.pptx")
                st.stop()
            file_list.append(chapter1_file)

            # 2. 第二章
            ch2_file = os.path.join(moban_dir, f"第二章{prefix}.pptx")
            if not os.path.exists(ch2_file):
                st.error(f"找不到第二章文件：{ch2_file}")
                st.stop()
            file_list.append(ch2_file)
            if domain == "具身智能":
                ch2_common = os.path.join(moban_dir, "第二章具身智能.pptx")
                if os.path.exists(ch2_common):
                    file_list.append(ch2_common)
                else:
                    st.warning("未找到第二章具身智能通用文件，跳过")

            # 3. 第三章
            ch3_file = os.path.join(moban_dir, f"第三章{prefix}.pptx")
            if not os.path.exists(ch3_file):
                st.error(f"找不到第三章文件：{ch3_file}")
                st.stop()
            file_list.append(ch3_file)
            if domain == "具身智能":
                ch3_common = os.path.join(moban_dir, "第三章具身智能.pptx")
                if os.path.exists(ch3_common):
                    file_list.append(ch3_common)
                else:
                    st.warning("未找到第三章具身智能通用文件，跳过")

            # 4. 第四章
            ch4_file = os.path.join(moban_dir, "第四章配套服务.pptx")
            if not os.path.exists(ch4_file):
                st.error("找不到第四章配套服务.pptx")
                st.stop()
            file_list.append(ch4_file)

            # 英伟达页（如果勾选）
            if need_nvidia == "是":
                nvidia_file = os.path.join(moban_dir, "英伟达.pptx")
                if os.path.exists(nvidia_file):
                    file_list.append(nvidia_file)
                else:
                    st.warning("未找到英伟达.pptx，已跳过")

            # 7.3 准备输出目录
            output_dir = os.path.join(os.getcwd(), "generated_ppt")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{school_name}_{domain}_{sub_domain or ''}方案.pptx")

            # 7.4 合并PPT
            try:
                merge_pptx_files(
                    file_paths=file_list,
                    output_path=output_path,
                    replace_cover_title=ppt_title,
                    replace_school_name=school_name
                )

                # 如果上传了图片，追加到末尾
                if uploaded_image:
                    prs = Presentation(output_path)
                    image_data = uploaded_image.read()
                    add_image_slide_to_end(prs, image_data)
                    prs.save(output_path)

                # 7.5 下载
                if os.path.exists(output_path):
                    with open(output_path, "rb") as f:
                        st.success("PPT生成成功！")
                        st.download_button(
                            label="下载PPT",
                            data=f,
                            file_name=f"{school_name}_{domain}_{sub_domain or ''}方案.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                else:
                    st.error("生成失败，未找到输出文件")

            except Exception as e:
                st.error(f"合并PPT时出错：{e}")
                st.exception(e)