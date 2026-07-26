# -*- coding: utf-8 -*-
"""turing_ppt_editor.py - 具身智能PPT模板原地修改 skill 脚本。
合并：正文处理(第3页公司介绍/第11页英伟达) + 封面标题/目录章节原地改。
保留模板排版对齐字体，只改文字。不依赖 cover_builder。"""
import argparse
import io
import os
import re
import shutil
import sys
from copy import deepcopy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.parts.image import Image

INTRO_MAX_CHARS = 500
SLIDE3_INDEX = 2
SLIDE11_INDEX = 10
INTRO_MARKER = "广东图灵智新"
SLIDE11_MARKER = "NVIDIA优选级解决方案合作伙伴"
NVIDIA_PAT = re.compile(r"\s*NVIDIA英伟达、?")
DEFAULT_OUT_DIR = r"D:\tulingppt"
CH1_FIXED = "图灵智新介绍"
MAX_CHAPTERS = 6
HIGHLIGHT_COLOR = "5DA0F9"

# ?? -> ??????????????(1-based)
THEME_SLIDE_MAP = {
    "具身智能": [1, 4, 7, 8, 9, 10, 11, 12, 13],
    "人工智能": [2, 5],
    "ai": [2, 5],
    "低空技术": [3, 6],
    "低空经济": [3, 6],
    "低空": [3, 6],
}
THEME_TO_CH2 = {
    "ai": "人工智能产业背景及国家政策",
    "人工智能": "人工智能产业背景及国家政策",
    "ai架构": "人工智能产业背景及国家政策",
    "具身智能": "具身智能机器人产业背景及国家政策",
    "低空技术": "低空技术产业背景及国家政策",
    "低空经济": "低空技术产业背景及国家政策",
    "低空": "低空技术产业背景及国家政策",
    "视频生成": "视频生成产业背景及国家政策",
    "ai视频生成": "视频生成产业背景及国家政策",
    "aigc视频": "视频生成产业背景及国家政策",
    "aigc": "AIGC产业背景及国家政策",
    "机器人": "机器人产业背景及国家政策",
    "大模型": "大模型产业背景及国家政策",
    "llm": "大模型产业背景及国家政策",
    "智能驾驶": "智能驾驶产业背景及国家政策",
    "自动驾驶": "自动驾驶产业背景及国家政策",
}


def walk(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(shp.shapes)


def find_intro_box(slide):
    for shp in walk(slide.shapes):
        if not shp.has_text_frame:
            continue
        t = shp.text_frame.text
        if INTRO_MARKER in t and len(t.replace("\n", "").replace("\r", "").replace("\x0b", "")) > 100:
            return shp
    return None


def find_slide11(prs):
    if len(prs.slides) <= SLIDE11_INDEX:
        raise ValueError(f"幻灯片总数 {len(prs.slides)} 不足 {SLIDE11_INDEX + 1}")
    slide = prs.slides[SLIDE11_INDEX]
    all_text = "".join(shp.text_frame.text for shp in walk(slide.shapes) if shp.has_text_frame)
    if SLIDE11_MARKER not in all_text and "英伟达" not in all_text:
        raise ValueError("第11页特征校验失败，可能不是本模板")
    return slide


def set_text_keep(shp, text):
    """原地改文字，保留首 run 格式，删除多余 run/段落。"""
    tf = shp.text_frame
    paras = list(tf.paragraphs)
    if not paras:
        shp.text_frame.text = text
        return
    p0 = paras[0]
    runs = list(p0.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
    else:
        p0.text = text
    for p in paras[1:]:
        p._p.getparent().remove(p._p)


def set_intro_text(intro_box, new_text):
    """替换公司介绍框文字。继承原段首缩进(前导空格)+过滤空行+复制首段pPr+deepcopy首run rPr。
    字数(含前导缩进空格的实际字符)≤500。建议3段。"""
    cleaned = new_text.replace("\r", "")
    tf = intro_box.text_frame
    paras = list(tf.paragraphs)
    first_p = paras[0]
    runs = list(first_p.runs)
    template_rPr = None
    if runs:
        rpr = runs[0]._r.find(qn("a:rPr"))
        if rpr is not None:
            template_rPr = deepcopy(rpr)
    template_pPr = deepcopy(first_p._p.find(qn("a:pPr")))
    orig_text = first_p.text
    leading = ""
    for ch in orig_text:
        if ch in " \t":
            leading += ch
        else:
            break
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)

    def make_run(p, line):
        r = p.add_run()
        r.text = line
        if template_rPr is not None:
            old = r._r.find(qn("a:rPr"))
            if old is not None:
                r._r.remove(old)
            r._r.insert(0, deepcopy(template_rPr))

    lines = [ln for ln in cleaned.split("\n") if ln.strip()]
    if not lines:
        lines = [""]
    lines = [leading + ln for ln in lines]
    n = len("".join(lines))
    if n > INTRO_MAX_CHARS:
        raise ValueError(f"公司介绍字数 {n} 超过上限 {INTRO_MAX_CHARS}")
    make_run(first_p, lines[0])
    for line in lines[1:]:
        p = tf.add_paragraph()
        if template_pPr is not None:
            old_pPr = p._p.find(qn("a:pPr"))
            if old_pPr is not None:
                p._p.remove(old_pPr)
            p._p.insert(0, deepcopy(template_pPr))
        make_run(p, line)
    return n


def _clean_nvidia_text(full):
    new = NVIDIA_PAT.sub("", full)
    new = new.replace("NVIDIA", "").replace("英伟达", "")
    new = re.sub(r"、{2,}", "、", new)
    new = re.sub(r"、\s+", "、", new)
    new = re.sub(r"\s{2,}", " ", new)
    new = re.sub(r"^\s*、|、\s*$", "", new)
    return new


def clean_nvidia_in_intro(intro_box):
    changed = False
    for p in intro_box.text_frame.paragraphs:
        runs = list(p.runs)
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        if "NVIDIA" not in full and "英伟达" not in full:
            continue
        new = _clean_nvidia_text(full)
        if new != full:
            runs[0].text = new
            for r in runs[1:]:
                r.text = ""
            changed = True
    return changed


def scrub_nvidia_all(prs):
    """清理全PPT所有文本框NVIDIA/英伟达字样。"""
    count = 0
    for slide in prs.slides:
        for shp in walk(slide.shapes):
            if not shp.has_text_frame:
                continue
            t = shp.text_frame.text
            if "NVIDIA" not in t and "英伟达" not in t:
                continue
            changed = False
            for p in shp.text_frame.paragraphs:
                runs = list(p.runs)
                if not runs:
                    continue
                full = "".join(r.text for r in runs)
                if "NVIDIA" not in full and "英伟达" not in full:
                    continue
                new = _clean_nvidia_text(full)
                if new != full:
                    runs[0].text = new
                    for r in runs[1:]:
                        r.text = ""
                    changed = True
            if changed:
                count += 1
    return count


def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.rId
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def _max_id(slide):
    m = 0
    for shp in walk(slide.shapes):
        cNvPr = shp._element.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None:
            try:
                m = max(m, int(cNvPr.get('id', '0')))
            except Exception:
                pass
    return m


def _clone_shape(shp, new_top, text, nid, spTree):
    new_el = deepcopy(shp._element)
    off = new_el.find('.//' + qn('a:off'))
    if off is not None:
        off.set('y', str(new_top))
    cNvPr = new_el.find('.//' + qn('p:cNvPr'))
    if cNvPr is not None:
        cNvPr.set('id', str(nid))
    ts = new_el.findall('.//' + qn('a:t'))
    if ts:
        ts[0].text = text
        for t in ts[1:]:
            t.text = ""
    spTree.append(new_el)


def set_cover_title(slide, title):
    """原地改第1页封面标题（找含'建议方案'或'解决方案'的短文本框）。"""
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text
        if ("建议方案" in txt or "解决方案" in txt) and len(txt) < 30:
            set_text_keep(shp, title)
            return True
    return False


def set_cover_school(slide, school):
    """在标题上方加学校名文本框（左上角区域，Arial 30pt粗体白色）。
    位置 (0.69, 2.2)in，标题#6在(0.69, 2.98)上方。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    for shp in list(slide.shapes):
        if shp.name == "校名" and shp.has_text_frame:
            shp._element.getparent().remove(shp._element)
    txBox = slide.shapes.add_textbox(Inches(0.69), Inches(2.2), Inches(11.14), Inches(0.6))
    txBox.name = "校名"
    tf = txBox.text_frame
    tf.text = school
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return True


def set_cover_badge(slide, badge_path):
    """左上角(0.5, 0.3)in插入校徽图片(1.5x0.7in)。可选，用户后续提供校徽时调用。
    默认不添加校徽，仅预留位置坐标。"""
    from pptx.util import Inches
    if badge_path and os.path.exists(badge_path):
        slide.shapes.add_picture(badge_path, Inches(0.5), Inches(0.3), Inches(1.5), Inches(0.7))
        return True
    return False


def set_toc(slide, chapters, current=1):
    """原地改第2页目录：按 chapters 数调整矩形对(克隆/删除)，改编号1..N和章节名。
    保留排版对齐和填充（标蓝位置不动）。返回章节数。"""
    rects = [shp for shp in slide.shapes if shp.name.startswith("Rectangle") and shp.has_text_frame]
    rects.sort(key=lambda s: (s.top, s.width))
    pairs = []
    used = set()
    for i, r in enumerate(rects):
        if i in used:
            continue
        for j in range(i + 1, len(rects)):
            if j in used:
                continue
            if abs(r.top - rects[j].top) < 100000:
                num = r if r.width < rects[j].width else rects[j]
                title = rects[j] if r.width < rects[j].width else r
                pairs.append((num, title))
                used.add(i)
                used.add(j)
                break
    pairs.sort(key=lambda p: p[0].top)
    N = len(pairs)
    M = len(chapters)
    row_h = int(1.103 * 914400)
    spTree = slide.shapes._spTree
    if M > N:
        maxid = _max_id(slide)
        for k in range(M - N):
            num_last, title_last = pairs[-1]
            new_y = num_last.top + row_h * (k + 1)
            _clone_shape(num_last, new_y, str(N + k + 1), maxid + 2 * k + 1, spTree)
            _clone_shape(title_last, new_y, chapters[N + k], maxid + 2 * k + 2, spTree)
    elif M < N:
        for num, title in pairs[M:]:
            num._element.getparent().remove(num._element)
            title._element.getparent().remove(title._element)
        pairs = pairs[:M]
    for i, (num, title) in enumerate(pairs[:M]):
        set_text_keep(num, str(i + 1))
        set_text_keep(title, chapters[i])
    return M


def map_theme_to_ch2(theme):
    """主题关键词 -> 第2章名「<主题>产业背景及国家政策」。内置常见主题映射，未命中则用主题名兜底。"""
    key = theme.strip().lower()
    if key in THEME_TO_CH2:
        return THEME_TO_CH2[key]
    if theme.strip() in THEME_TO_CH2:
        return THEME_TO_CH2[theme.strip()]
    return f"{theme.strip()}产业背景及国家政策"


def build_toc_from_theme(theme, toc_tail):
    """目录规则：ch1固定「图灵智新介绍」+ ch2「<主题>产业背景及国家政策」+ toc_tail(第3-6章自由排版)。总章数≤6。"""
    chapters = [CH1_FIXED, map_theme_to_ch2(theme)]
    if toc_tail:
        chapters.extend(c.strip() for c in toc_tail.split("|") if c.strip())
    if len(chapters) > MAX_CHAPTERS:
        raise ValueError(f"目录章数 {len(chapters)} 超过上限 {MAX_CHAPTERS}（ch1图灵智新介绍+ch2主题背景+第3-6章）")
    return chapters


def resolve_chapters(args):
    """解析目录章节：--set-toc 优先(全手动)，否则 --theme 自动建(ch1固定+ch2主题+tail)，否则 None。"""
    if args.set_toc:
        return args.set_toc.split("|")
    if args.theme:
        return build_toc_from_theme(args.theme, args.toc_tail)
    return None


def resolve_out_path(out, pptx):
    """输出路径改为当前目录下的 generated_ppt 文件夹"""
    out_dir = os.path.join(os.getcwd(), "generated_ppt")
    os.makedirs(out_dir, exist_ok=True)
    if out:
        basename = os.path.basename(out)
        if not basename:
            basename = os.path.basename(pptx)
    else:
        basename = os.path.basename(pptx)
    return os.path.join(out_dir, basename)


# === NEW FUNCTIONS ===

def duplicate_slide(prs, index):
    """Clone a slide within the same presentation. Returns the new slide (appended at end)."""
    source = prs.slides[index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    for shp in source.shapes:
        el = deepcopy(shp._element)
        new_slide.shapes._spTree.append(el)
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype or "tags" in rel.reltype:
            continue
        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
    return new_slide


def _find_toc_pairs(slide):
    """Find Rectangle pairs (number + title) on a TOC slide, sorted by top."""
    rects = [shp for shp in slide.shapes if shp.name.startswith("Rectangle") and shp.has_text_frame]
    rects.sort(key=lambda s: (s.top, s.width))
    pairs = []
    used = set()
    for i, r in enumerate(rects):
        if i in used:
            continue
        for j in range(i + 1, len(rects)):
            if j in used:
                continue
            if abs(r.top - rects[j].top) < 100000:
                num = r if r.width < rects[j].width else rects[j]
                title = rects[j] if r.width < rects[j].width else r
                pairs.append((num, title))
                used.add(i)
                used.add(j)
                break
    pairs.sort(key=lambda p: p[0].top)
    return pairs


def set_toc_highlight(slide, current):
    """Change which chapter is highlighted blue in the TOC.
    Highlighted (current): blue fill (srgbClr HIGHLIGHT_COLOR) + white text (schemeClr bg1)
    Non-highlighted: white fill (schemeClr bg1) + inherited dark text (remove explicit color)
    """
    pairs = _find_toc_pairs(slide)
    for i, (num, title) in enumerate(pairs):
        is_current = (i + 1 == current)
        for shp in (num, title):
            # --- Change fill ---
            spPr = shp._element.find(qn("p:spPr"))
            if spPr is not None:
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    for child in list(solidFill):
                        solidFill.remove(child)
                    if is_current:
                        clr = etree.SubElement(solidFill, qn("a:srgbClr"))
                        clr.set("val", HIGHLIGHT_COLOR)
                    else:
                        clr = etree.SubElement(solidFill, qn("a:schemeClr"))
                        clr.set("val", "bg1")

            # --- Change text color ---
            for paragraph in shp.text_frame.paragraphs:
                for run in paragraph.runs:
                    rPr = run._r.find(qn("a:rPr"))
                    if is_current:
                        # Set white text (schemeClr bg1 = white)
                        if rPr is None:
                            rPr = etree.Element(qn("a:rPr"))
                            run._r.insert(0, rPr)
                        existing_sf = rPr.find(qn("a:solidFill"))
                        if existing_sf is not None:
                            rPr.remove(existing_sf)
                        sf = etree.SubElement(rPr, qn("a:solidFill"))
                        clr = etree.SubElement(sf, qn("a:schemeClr"))
                        clr.set("val", "bg1")
                    else:
                        # Remove explicit text color (inherit dark from style)
                        if rPr is not None:
                            existing_sf = rPr.find(qn("a:solidFill"))
                            if existing_sf is not None:
                                rPr.remove(existing_sf)



def get_theme_colors(prs):
    """Extract theme color scheme from presentation. Returns dict of color_name -> RGB hex."""
    from pptx.oxml.ns import qn
    from lxml import etree
    master = prs.slide_masters[0]
    theme_part = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:
        return {}
    theme_el = etree.fromstring(theme_part.blob)
    clrScheme = theme_el.find('.//' + qn('a:clrScheme'))
    if clrScheme is None:
        return {}
    colors = {}
    for child in clrScheme:
        name = child.tag.split('}')[-1]
        srgb = child.find(qn('a:srgbClr'))
        sysClr = child.find(qn('a:sysClr'))
        if srgb is not None:
            colors[name] = srgb.get('val')
        elif sysClr is not None:
            colors[name] = sysClr.get('lastClr', sysClr.get('val', ''))
    return colors


def get_effective_clr_map(prs, slide):
    """Get effective color map for a slide (considering clrMapOvr)."""
    from pptx.oxml.ns import qn
    clrMapOvr = slide._element.find(qn('p:clrMapOvr'))
    if clrMapOvr is not None:
        override = clrMapOvr.find(qn('a:overrideClrMapping'))
        if override is not None:
            return dict(override.attrib)
    master = prs.slide_masters[0]
    clrMap = master.element.find('.//' + qn('p:clrMap'))
    if clrMap is not None:
        return dict(clrMap.attrib)
    return {}


def convert_scheme_colors(xml_element, theme_colors, clr_map):
    """Convert all schemeClr elements to srgbClr using theme_colors and clr_map.
    Preserves child elements (lumMod, tint, shade, alpha, etc.)."""
    from lxml import etree
    A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    scheme_elements = list(xml_element.iter(A_NS + 'schemeClr'))
    for schemeClr in scheme_elements:
        val = schemeClr.get('val')
        if val in clr_map:
            theme_name = clr_map[val]
        else:
            theme_name = val
        rgb = theme_colors.get(theme_name)
        if rgb:
            srgbClr = etree.Element(A_NS + 'srgbClr')
            srgbClr.set('val', rgb)
            for child in list(schemeClr):
                srgbClr.append(child)
            parent = schemeClr.getparent()
            idx = list(parent).index(schemeClr)
            parent.remove(schemeClr)
            parent.insert(idx, srgbClr)

def copy_external_slide(prs, source_slide, source_prs=None):
    """Copy a slide from another presentation, preserving background and colors.
    Converts theme colors (schemeClr) to explicit RGB (srgbClr) to prevent
    color shifts when the target presentation has a different theme."""
    from io import BytesIO
    from copy import deepcopy
    from pptx.oxml.ns import qn

    # Extract theme colors and color map from source presentation
    if source_prs is not None:
        theme_colors = get_theme_colors(source_prs)
        clr_map = get_effective_clr_map(source_prs, source_slide)
    else:
        theme_colors = {}
        clr_map = {}

    try:
        layout = prs.slide_layouts[6]
    except IndexError:
        layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(layout)

    # Copy background from source slide
    source_cSld = source_slide._element.find(qn('p:cSld'))
    new_cSld = new_slide._element.find(qn('p:cSld'))
    if source_cSld is not None and new_cSld is not None:
        source_bg = source_cSld.find(qn('p:bg'))
        if source_bg is not None:
            existing_bg = new_cSld.find(qn('p:bg'))
            if existing_bg is not None:
                new_cSld.remove(existing_bg)
            bg_copy = deepcopy(source_bg)
            convert_scheme_colors(bg_copy, theme_colors, clr_map)
            new_cSld.insert(0, bg_copy)

    # Copy color map override
    source_clrMap = source_slide._element.find(qn('p:clrMapOvr'))
    if source_clrMap is not None:
        existing = new_slide._element.find(qn('p:clrMapOvr'))
        if existing is not None:
            new_slide._element.remove(existing)
        new_slide._element.append(deepcopy(source_clrMap))

    # Remove default shapes
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    # Build rId mapping for images
    rid_map = {}
    for rel in source_slide.part.rels.values():
        if "image" not in rel.reltype and not rel.is_external:
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            rid_map[rel.rId] = new_rid
        elif "image" in rel.reltype:
            blob = rel.target_part.blob
            ct = rel.target_part.content_type
            new_rid = None
            # 1) Try add_picture directly (PNG/JPEG/GIF/BMP/TIFF/WMF)
            try:
                stream = BytesIO(blob)
                pic = new_slide.shapes.add_picture(stream, 0, 0)
                blip = pic._element.find('.//' + qn('a:blip'))
                if blip is not None:
                    new_rid = blip.get(qn('r:embed'))
                pic._element.getparent().remove(pic._element)
            except Exception:
                pass
            # 2) WebP: convert to PNG via Pillow
            if new_rid is None and ct == "image/webp":
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(BytesIO(blob))
                    png_stream = BytesIO()
                    img.save(png_stream, format='PNG')
                    png_stream.seek(0)
                    pic = new_slide.shapes.add_picture(png_stream, 0, 0)
                    blip = pic._element.find('.//' + qn('a:blip'))
                    if blip is not None:
                        new_rid = blip.get(qn('r:embed'))
                    pic._element.getparent().remove(pic._element)
                except Exception as e:
                    print(f"  WARN: WebP conversion failed: {e}")
            # 3) SVG / other: add as raw part (PowerPoint supports SVG natively)
            if new_rid is None:
                try:
                    ext = ct.split('/')[-1].replace('+xml', '')
                    package = prs.part.package
                    partname = package.next_partname("/ppt/media/image%d." + ext)

                    class _RawImagePart(Part):
                        def __init__(self, partname, content_type, blob_data, package=None):
                            super().__init__(partname, content_type, package)
                            self._blob_data = blob_data
                        @property
                        def blob(self):
                            return self._blob_data

                    raw_part = _RawImagePart(partname, ct, blob, package)
                    new_rid = new_slide.part.relate_to(raw_part, RT.IMAGE)
                except Exception as e:
                    print(f"  WARN: raw image part failed: {e}")
            if new_rid is not None:
                rid_map[rel.rId] = new_rid

    # Copy shapes with updated rIds and converted colors
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for shp in source_slide.shapes:
        el = deepcopy(shp._element)
        for elem in el.iter():
            for attr_name in ["{%s}embed" % R_NS, "{%s}link" % R_NS, "{%s}id" % R_NS]:
                old_val = elem.get(attr_name)
                if old_val and old_val in rid_map:
                    elem.set(attr_name, rid_map[old_val])
        convert_scheme_colors(el, theme_colors, clr_map)
        new_slide.shapes._spTree.append(el)
    return new_slide


def append_theme_content(prs, merge_path, theme, extra_path=None):
    """Append theme-specific slides from merge source and optional extra pptx.
    Preserves source slide backgrounds and converts theme colors to explicit RGB."""
    count = 0
    key = theme.strip().lower()
    if key not in THEME_SLIDE_MAP:
        key = theme.strip()
    slide_indices = THEME_SLIDE_MAP.get(key, [])
    if merge_path and os.path.exists(merge_path) and slide_indices:
        source_prs = Presentation(merge_path)
        for idx in slide_indices:
            if idx <= len(source_prs.slides):
                copy_external_slide(prs, source_prs.slides[idx - 1], source_prs)
                count += 1
    if extra_path and os.path.exists(extra_path):
        extra_prs = Presentation(extra_path)
        for slide in extra_prs.slides:
            copy_external_slide(prs, slide, extra_prs)
            count += 1
    return count

def insert_toc2(prs, current=2):
    """Insert a duplicate of slide 2 (TOC) at the end, with specified chapter highlighted."""
    new_slide = duplicate_slide(prs, 1)
    set_toc_highlight(new_slide, current)
    return new_slide

def build_full(pptx, title, chapters, current, intro, no_nvidia, school, badge, out,
               insert_toc2_flag=False, merge_source=None, theme_select=None,
               extra_pptx=None, chapters_filter=None):
    """一键原地改：封面标题+校名+校徽 + 目录章节 + 第3页介绍 + 删英伟达 + 清NVIDIA +
    可选: 插入TOC2页(标蓝第二章) + 合并主题内容 + 追加外部PPTX。"""
    prs = Presentation(pptx)
    if title:
        if set_cover_title(prs.slides[0], title):
            print(f"✅ 第1页标题: {title}")
        else:
            print("⚠️ 第1页未找到标题框")
    if school:
        set_cover_school(prs.slides[0], school)
        print(f"✅ 第1页校名: {school} (标题上方)")
    if badge:
        if set_cover_badge(prs.slides[0], badge):
            print(f"✅ 第1页校徽: {badge} (左上角)")
        else:
            print("⚠️ 校徽文件不存在")
    if chapters:
        m = set_toc(prs.slides[1], chapters, current)
        print(f"✅ 第2页目录: {m}章")
    if intro:
        intro_box = find_intro_box(prs.slides[SLIDE3_INDEX])
        if intro_box:
            n = set_intro_text(intro_box, intro)
            print(f"✅ 第3页公司介绍: {n}字")

    # 确定是否包含第二章
    include_ch2 = chapters_filter is None or "2" in chapters_filter

    # 先插入TOC2和主题内容（在删NVIDIA之前，避免partname冲突）
    if insert_toc2_flag and include_ch2:
        insert_toc2(prs, current=2)
        print("✅ 插入TOC2页(标蓝第二章)")
    if include_ch2 and theme_select and (merge_source or extra_pptx):
        n = append_theme_content(prs, merge_source, theme_select, extra_pptx)
        print(f"✅ 追加主题内容: {n}页")

    # 后删NVIDIA（避免partname重复冲突）
    if no_nvidia:
        intro_box = find_intro_box(prs.slides[SLIDE3_INDEX])
        if intro_box:
            clean_nvidia_in_intro(intro_box)
        try:
            find_slide11(prs)
            delete_slide(prs, SLIDE11_INDEX)
            print("✅ 第11页(英伟达)已删")
        except Exception:
            pass
        scrub_nvidia_all(prs)
        print("✅ 清NVIDIA字样")

    out_path = resolve_out_path(out, pptx)
    prs.save(out_path)
    saved = Presentation(out_path)
    print(f"✅ 已保存: {out_path}（{len(saved.slides)}页）")


def main():
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")
    ap = argparse.ArgumentParser(description="PPT模板原地修改器（turing-ppt-editor skill）")
    ap.add_argument("--pptx", required=True, help="目标 PPTX 路径")
    ap.add_argument("--full", action="store_true", help="一键原地改: 标题+目录+介绍+删英伟达+可选合并")
    ap.add_argument("--set-title", help="原地改第1页封面标题")
    ap.add_argument("--set-school", help="封面加学校名(标题上方,Arial 30pt粗体白)")
    ap.add_argument("--set-badge", help="封面左上角插入校徽图片(可选,默认预留位置不加)")
    ap.add_argument("--set-toc", help="原地改第2页目录,|分隔章节(动态章节数,克隆/删除矩形)")
    ap.add_argument("--theme", help="主题(ai/具身智能/低空技术/视频生成...),自动建目录+主题幻灯片选择")
    ap.add_argument("--toc-tail", help="目录第3-6章,|分隔(自由排版,与--theme配合,总章数≤6)")
    ap.add_argument("--current", type=int, default=1, help="标蓝章节序号(从1,默认第1章)")
    ap.add_argument("--set-intro", help="改第3页公司介绍(≤500字)")
    ap.add_argument("--intro-file", help="从文件读公司介绍(≤500字,支持多行)")
    ap.add_argument("--no-nvidia", action="store_true", help="删第11页+清第3页英伟达字样")
    ap.add_argument("--scrub-nvidia", action="store_true", help="清全PPT所有NVIDIA/英伟达字样")
    ap.add_argument("--remove-slide11", action="store_true", help="仅删第11页")
    ap.add_argument("--insert-toc2", action="store_true", help="第一章后插入TOC2页(标蓝第二章)")
    ap.add_argument("--merge-source", help="合并源PPTX(固定用D:\moban\merge_3themes.pptx),按主题选择幻灯片")
    ap.add_argument("--theme-select", help="主题选择(具身智能/人工智能/低空技术),从合并源选对应幻灯片")
    ap.add_argument("--extra-pptx", help="第3章内容PPTX(从D:\moban\选对应文件,见SKILL.md映射表)")
    ap.add_argument("--chapters", help="包含的章节(如'1,2'),默认全部")
    ap.add_argument("--out", help="输出文件名(目录固定D:/tulingppt,仅取文件名部分)")
    args = ap.parse_args()

    if not (args.full or args.set_title or args.set_school or args.set_badge or args.set_toc or args.set_intro is not None
            or args.intro_file or args.no_nvidia or args.scrub_nvidia or args.remove_slide11 or args.theme
            or args.insert_toc2 or args.merge_source):
        ap.error("至少指定一个操作: --full / --set-title / --set-school / --set-badge / --set-toc / --set-intro / --intro-file / --no-nvidia / --scrub-nvidia / --remove-slide11 / --theme / --insert-toc2 / --merge-source")
    if not os.path.exists(args.pptx):
        sys.exit(f"❌ 文件不存在: {args.pptx}")

    if args.full:
        title = args.set_title
        chapters = resolve_chapters(args)
        intro = args.set_intro
        if args.intro_file:
            with open(args.intro_file, encoding="utf-8-sig") as f:
                intro = f.read()
        theme_select = args.theme_select or args.theme
        build_full(args.pptx, title, chapters, args.current, intro, args.no_nvidia,
                   args.set_school, args.set_badge, args.out,
                   insert_toc2_flag=args.insert_toc2, merge_source=args.merge_source,
                   theme_select=theme_select, extra_pptx=args.extra_pptx,
                   chapters_filter=args.chapters)
        return

    prs = Presentation(args.pptx)
    actions = []
    if args.set_title:
        if set_cover_title(prs.slides[0], args.set_title):
            actions.append(f"第1页标题改: {args.set_title}")
    if args.set_school:
        set_cover_school(prs.slides[0], args.set_school)
        actions.append(f"第1页校名: {args.set_school}")
    if args.set_badge:
        if set_cover_badge(prs.slides[0], args.set_badge):
            actions.append(f"第1页校徽: {args.set_badge}")
    chapters = resolve_chapters(args)
    if chapters:
        m = set_toc(prs.slides[1], chapters, args.current)
        actions.append(f"第2页目录: {m}章")
    if args.set_intro is not None or args.intro_file:
        intro = args.set_intro
        if args.intro_file:
            with open(args.intro_file, encoding="utf-8-sig") as f:
                intro = f.read()
        intro_box = find_intro_box(prs.slides[SLIDE3_INDEX])
        if intro_box:
            n = set_intro_text(intro_box, intro)
            actions.append(f"第3页介绍: {n}字")
    if args.no_nvidia:
        intro_box = find_intro_box(prs.slides[SLIDE3_INDEX])
        if intro_box:
            clean_nvidia_in_intro(intro_box)
        try:
            find_slide11(prs)
            delete_slide(prs, SLIDE11_INDEX)
            actions.append("第11页英伟达已删")
        except Exception:
            pass
        scrub_nvidia_all(prs)
        actions.append("清NVIDIA字样")
    if args.scrub_nvidia and not args.no_nvidia:
        c = scrub_nvidia_all(prs)
        actions.append(f"清NVIDIA: {c}框")
    if args.remove_slide11 and not args.no_nvidia:
        try:
            find_slide11(prs)
            delete_slide(prs, SLIDE11_INDEX)
            actions.append("第11页已删")
        except Exception:
            pass
    if args.insert_toc2:
        insert_toc2(prs, current=2)
        actions.append("插入TOC2页(标蓝第二章)")
    if args.merge_source and args.theme_select:
        n = append_theme_content(prs, args.merge_source, args.theme_select, args.extra_pptx)
        actions.append(f"追加主题内容: {n}页")
    out_path = resolve_out_path(args.out, args.pptx)
    prs.save(out_path)
    actions.append(f"保存: {out_path}")
    print("\n".join("✅ " + a for a in actions))


if __name__ == "__main__":
    try:
        main()
    except (ValueError, FileNotFoundError) as e:
        sys.exit(f"[ERROR] {e}")
    except Exception as e:
        sys.exit(f"[ERROR] {type(e).__name__}: {e}")
